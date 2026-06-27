"""
PPTX 面试PPT生成服务
根据 AI 生成的 PPT 内容创建 PowerPoint 文件
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os


def generate_ppt(ppt_content: dict, output_path: str):
    """根据内容结构生成 PPTX 文件"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)

    slides_data = ppt_content.get("slides", [])
    title = ppt_content.get("title", "面试演示")
    subtitle = ppt_content.get("subtitle", "")

    # 颜色方案
    PRIMARY = RGBColor(0x15, 0x70, 0xEF)  # 主色
    DARK = RGBColor(0x10, 0x18, 0x28)      # 深色文字
    SUB = RGBColor(0x47, 0x54, 0x67)       # 副色文字
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFA)
    ACCENT_BG = RGBColor(0xEF, 0xF4, 0xFF)

    # ── 封面页 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    # 背景色
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY

    # 标题
    left = Cm(3)
    top = Cm(4)
    width = Cm(34)
    height = Cm(3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # 副标题
    if subtitle:
        top2 = Cm(6.5)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Cm(2))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xEE)
        p2.alignment = PP_ALIGN.LEFT

    # ── 内容页 ──
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 白色背景
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # 顶部色条
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Cm(0), Cm(0), prs.slide_width, Cm(0.35)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()

        # 标题
        slide_title = slide_data.get("title", f"第{i+1}页")
        txBox = slide.shapes.add_textbox(Cm(2), Cm(1.2), Cm(30), Cm(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK

        # 内容
        content = slide_data.get("content", [])
        content_top = Cm(3.5)
        txBox2 = slide.shapes.add_textbox(Cm(2.5), content_top, Cm(29), Cm(12))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for j, item in enumerate(content):
            if j == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()

            p.text = f"•  {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = SUB
            p.space_after = Pt(12)
            p.alignment = PP_ALIGN.LEFT

        # 页脚
        txBox3 = slide.shapes.add_textbox(Cm(2), Cm(17.5), Cm(30), Cm(1))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f"{i+1} / {len(slides_data)}"
        p3.font.size = Pt(10)
        p3.font.color.rgb = RGBColor(0x98, 0xA2, 0xB3)
        p3.alignment = PP_ALIGN.RIGHT

        # 备注
        notes_text = slide_data.get("notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    prs.save(output_path)
    return output_path